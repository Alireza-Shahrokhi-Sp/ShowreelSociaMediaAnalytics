from pptx import Presentation
path = r"D:\Polythecninco di Milano\AFB_Lab\Ali\outputs\modeling\Instagram_Modeling_Report.pptx"
prs = Presentation(path)
for i, slide in enumerate(prs.slides, 1):
    texts = []
    for shape in slide.shapes:
        if shape.has_text_frame:
            t = shape.text_frame.text.strip()
            if t:
                texts.append(t[:120])
        elif shape.shape_type == 19:  # TABLE
            for row in shape.table.rows:
                for cell in row.cells:
                    t = cell.text_frame.text.strip()
                    if t:
                        texts.append(t[:60])
    print(f"=== SLIDE {i} ===")
    for t in texts[:15]:
        print(f"  {repr(t)}".encode('ascii', errors='replace').decode('ascii'))
    print()
print(f"TOTAL SLIDES: {len(prs.slides)}")
